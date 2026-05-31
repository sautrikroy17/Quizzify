import os
from PIL import Image as PILImage
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# ─── BRAND COLORS ────────────────────────────────────────────────────────────
ACCENT1     = RGBColor(79, 172, 254)   # #4facfe  (primary blue)
ACCENT2     = RGBColor(0,  242, 254)   # #00f2fe  (secondary cyan)
WHITE       = RGBColor(255, 255, 255)
LIGHT_TEXT  = RGBColor(200, 214, 229)  # soft white-blue body text
MUTED       = RGBColor(100, 116, 139)  # muted gray

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def add_bg(slide):
    """Add the premium glowing background to the slide."""
    slide.shapes.add_picture("premium_bg.png", 0, 0, SLIDE_W, SLIDE_H)

def add_glass_card(slide, l, t, w, h):
    """Add the semi-transparent glassmorphism card."""
    slide.shapes.add_picture("glass_card.png", l, t, w, h)

def add_slide_number(slide, number):
    tf = slide.shapes.add_textbox(SLIDE_W - Inches(1.0), Inches(0.4), Inches(0.6), Inches(0.6))
    p = tf.text_frame.paragraphs[0]
    p.text = str(number)
    p.alignment = PP_ALIGN.RIGHT
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT2
    p.font.bold = True

def add_title(slide, text, subtitle=None):
    if subtitle:
        tf_sub = slide.shapes.add_textbox(Inches(0.8), Inches(0.55), Inches(11.0), Inches(0.4))
        p_sub = tf_sub.text_frame.paragraphs[0]
        p_sub.text = subtitle.upper()
        p_sub.font.size = Pt(14)
        p_sub.font.bold = True
        p_sub.font.color.rgb = ACCENT2

    tf = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.0), Inches(0.8))
    p = tf.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE

def fit_image(slide, img_path, area_l, area_t, area_w, area_h, padding=Inches(0.2)):
    al = area_l + padding
    at = area_t + padding
    aw = area_w - 2*padding
    ah = area_h - 2*padding

    try:
        im = PILImage.open(img_path)
        iw, ih = im.size
        ratio = iw / ih
        area_ratio = aw / ah

        if ratio > area_ratio:
            fit_w = aw
            fit_h = aw / ratio
        else:
            fit_h = ah
            fit_w = ah * ratio

        cl = al + (aw - fit_w) / 2
        ct = at + (ah - fit_h) / 2

        slide.shapes.add_picture(img_path, cl, ct, width=fit_w, height=fit_h)
    except Exception as e:
        print(f"WARN: Could not add image {img_path}: {e}")

# ─── SLIDE 1: TITLE ──────────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(BLANK)
add_bg(slide1)

# Main Title Container
add_glass_card(slide1, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.1))

tf = slide1.shapes.add_textbox(Inches(1.2), Inches(1.5), Inches(10.0), Inches(0.5))
p = tf.text_frame.paragraphs[0]
p.text = "21CSC101P / OODP"
p.font.size = Pt(16)
p.font.color.rgb = ACCENT2
p.font.bold = True

tf2 = slide1.shapes.add_textbox(Inches(1.2), Inches(2.0), Inches(10.0), Inches(1.2))
p2 = tf2.text_frame.paragraphs[0]
p2.text = "QUIZZIFY"
p2.font.size = Pt(72)
p2.font.color.rgb = WHITE
p2.font.bold = True

tf3 = slide1.shapes.add_textbox(Inches(1.2), Inches(3.2), Inches(10.0), Inches(0.5))
p3 = tf3.text_frame.paragraphs[0]
p3.text = "Project Development Presentation"
p3.font.size = Pt(20)
p3.font.color.rgb = LIGHT_TEXT

# Team members at bottom of card
tf5 = slide1.shapes.add_textbox(Inches(1.2), Inches(4.5), Inches(10.0), Inches(1.5))
tf5.text_frame.word_wrap = True
members = [
    ("TEAM MEMBERS", Pt(12), ACCENT2, True),
    ("Sautrik Roy [RA251103010052]     SHIVAANI M [RA2511003010066]", Pt(14), WHITE, False),
    ("Ashwin Kumar N [RA2511003010033]     Ayush Agarwal [RA2511003010012]", Pt(14), WHITE, False),
]
for i, (text, size, color, bold) in enumerate(members):
    p = tf5.text_frame.paragraphs[0] if i == 0 else tf5.text_frame.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.space_before = Pt(4)

# ─── SLIDE 2: INTRODUCTION ───────────────────────────────────────────────────
slide2 = prs.slides.add_slide(BLANK)
add_bg(slide2)
add_slide_number(slide2, 2)
add_title(slide2, "Platform Overview", "Introduction")

card_specs = [
    ("What is Quizzify?",
     "An intelligent, AI-powered Quiz Management System that auto-generates assessments from uploaded documents (PDF / PPTX)."),
    ("Core Workflow",
     "User uploads a document → Content is extracted → Quiz is forged automatically → Student takes the quiz and gets instant results."),
    ("Key Features",
     "Real-time leaderboards · Result analytics · AI-powered Forge engine · Secure, auth-protected accounts."),
    ("Technology Stack",
     "Backend Integration: C++ Core Architecture\nFrontend Interface: Javascript & Interactive UI Components"),
]

for i, (heading, body) in enumerate(card_specs):
    col = i % 2
    row = i // 2
    cl = Inches(0.8) + col * Inches(6.0)
    ct = Inches(2.0) + row * Inches(2.5)
    cw = Inches(5.6)
    ch = Inches(2.2)
    
    add_glass_card(slide2, cl, ct, cw, ch)
    
    th = slide2.shapes.add_textbox(cl + Inches(0.3), ct + Inches(0.2), cw - Inches(0.6), Inches(0.4))
    tp = th.text_frame.paragraphs[0]
    tp.text = heading
    tp.font.size = Pt(18)
    tp.font.color.rgb = ACCENT1
    tp.font.bold = True
    
    tb = slide2.shapes.add_textbox(cl + Inches(0.3), ct + Inches(0.7), cw - Inches(0.6), Inches(1.3))
    tb.text_frame.word_wrap = True
    tbp = tb.text_frame.paragraphs[0]
    tbp.text = body
    tbp.font.size = Pt(15)
    tbp.font.color.rgb = LIGHT_TEXT

# ─── SLIDE 3: OODP FEATURES ──────────────────────────────────────────────────
slide3 = prs.slides.add_slide(BLANK)
add_bg(slide3)
add_slide_number(slide3, 3)
add_title(slide3, "Core Concepts Applied", "OODP Features Utilized")

oodp_items = [
    ("Encapsulation", "User and QuizResult classes bundle related attributes with specialized methods, keeping internal states hidden from external interference."),
    ("Abstraction", "Complex backend AI generation and C++ database operations are hidden behind simple frontend Javascript endpoints."),
    ("Inheritance", "The system leverages hierarchical data models where specialized components extend base identity traits."),
    ("Polymorphism", "Methods are designed to dynamically handle varying data types, enabling seamless processing of both PDF and PPTX formats uniformly."),
]

for i, (concept, desc) in enumerate(oodp_items):
    ct = Inches(2.0) + i * Inches(1.25)
    add_glass_card(slide3, Inches(0.8), ct, Inches(11.733), Inches(1.0))
    
    th = slide3.shapes.add_textbox(Inches(1.1), ct + Inches(0.3), Inches(2.5), Inches(0.4))
    tp = th.text_frame.paragraphs[0]
    tp.text = concept
    tp.font.size = Pt(20)
    tp.font.color.rgb = ACCENT2
    tp.font.bold = True

    tb = slide3.shapes.add_textbox(Inches(3.8), ct + Inches(0.2), Inches(8.4), Inches(0.6))
    tb.text_frame.word_wrap = True
    tbp = tb.text_frame.paragraphs[0]
    tbp.text = desc
    tbp.font.size = Pt(16)
    tbp.font.color.rgb = LIGHT_TEXT

# ─── SLIDES 4-7: UML DIAGRAMS ────────────────────────────────────────────────
uml_slides = [
    ("Use Case Diagram",   "System Interactions",  "use_case.png"),
    ("Sequence Diagram",   "Data Flow Map",        "sequence.png"),
    ("Class Diagram",      "Entity Architecture",  "class.png"),
    ("Collaboration Diagram","Object Messaging",   "collaboration.png"),
]

for idx, (title, label, img) in enumerate(uml_slides):
    slide_num = 4 + idx
    s = prs.slides.add_slide(BLANK)
    add_bg(s)
    add_slide_number(s, slide_num)
    add_title(s, title, f"UML Diagram — {label}")

    img_area_l = Inches(0.8)
    img_area_t = Inches(1.8)
    img_area_w = Inches(11.733)
    img_area_h = Inches(5.0)
    
    add_glass_card(s, img_area_l, img_area_t, img_area_w, img_area_h)
    fit_image(s, img, img_area_l, img_area_t, img_area_w, img_area_h)

# ─── SLIDE 8: LANDING PAGE ───────────────────────────────────────────────────
s8 = prs.slides.add_slide(BLANK)
add_bg(s8)
add_slide_number(s8, 8)
add_title(s8, "Landing Page Interface", "User Interface")

add_glass_card(s8, Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))
fit_image(s8, "Landing page.jpeg", Inches(0.8), Inches(1.8), Inches(11.733), Inches(5.0))

# ─── SLIDE 9: INPUT / OUTPUT ─────────────────────────────────────────────────
s9 = prs.slides.add_slide(BLANK)
add_bg(s9)
add_slide_number(s9, 9)
add_title(s9, "Input & Output Screens", "User Interface")

# Left image
add_glass_card(s9, Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))
fit_image(s9, "screenshot_forge_upload.png", Inches(0.8), Inches(1.8), Inches(5.7), Inches(5.0))

# Right image
add_glass_card(s9, Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))
fit_image(s9, "screenshot_active_quiz.png", Inches(6.8), Inches(1.8), Inches(5.7), Inches(5.0))

# Labels
tf_l = s9.shapes.add_textbox(Inches(0.8), Inches(6.8), Inches(5.7), Inches(0.4))
p_l = tf_l.text_frame.paragraphs[0]
p_l.text = "INPUT: Document Upload"
p_l.font.size = Pt(14)
p_l.font.color.rgb = ACCENT2
p_l.alignment = PP_ALIGN.CENTER

tf_r = s9.shapes.add_textbox(Inches(6.8), Inches(6.8), Inches(5.7), Inches(0.4))
p_r = tf_r.text_frame.paragraphs[0]
p_r.text = "OUTPUT: Active Quiz Generation"
p_r.font.size = Pt(14)
p_r.font.color.rgb = ACCENT2
p_r.alignment = PP_ALIGN.CENTER

# ─── SLIDE 10: CONCLUSION ────────────────────────────────────────────────────
s10 = prs.slides.add_slide(BLANK)
add_bg(s10)
add_slide_number(s10, 10)
add_title(s10, "Summary & Future Scope", "Conclusion")

add_glass_card(s10, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.5))

conclusion_items = [
    "Quizzify demonstrates the robust capabilities of mapping Object-Oriented paradigms into real-world applications.",
    "The C++ backend guarantees high performance, while Javascript on the frontend ensures a seamless user experience.",
    "AI automation dramatically decreases time spent crafting educational assessments.",
    "Future scope involves further integration of specialized models and real-time collaboration features."
]

for i, text in enumerate(conclusion_items):
    ct = Inches(2.2) + i * Inches(0.9)
    
    # Bullet
    tf_b = s10.shapes.add_textbox(Inches(1.2), ct, Inches(0.5), Inches(0.5))
    p_b = tf_b.text_frame.paragraphs[0]
    p_b.text = "✦"
    p_b.font.size = Pt(20)
    p_b.font.color.rgb = ACCENT1

    tb = s10.shapes.add_textbox(Inches(1.8), ct + Inches(0.05), Inches(10.0), Inches(0.8))
    tb.text_frame.word_wrap = True
    tp_ = tb.text_frame.paragraphs[0]
    tp_.text = text
    tp_.font.size = Pt(18)
    tp_.font.color.rgb = LIGHT_TEXT

# (Thank you text moved to slide 11)

# ─── SLIDE 11: THANK YOU ─────────────────────────────────────────────────────
s11 = prs.slides.add_slide(BLANK)
add_bg(s11)

# Center glass card
add_glass_card(s11, Inches(2.666), Inches(2.0), Inches(8.0), Inches(3.5))

tf_ty_main = s11.shapes.add_textbox(Inches(2.666), Inches(2.7), Inches(8.0), Inches(1.0))
p_ty_main = tf_ty_main.text_frame.paragraphs[0]
p_ty_main.text = "THANK YOU"
p_ty_main.alignment = PP_ALIGN.CENTER
p_ty_main.font.size = Pt(64)
p_ty_main.font.bold = True
p_ty_main.font.color.rgb = WHITE

tf_ty_sub = s11.shapes.add_textbox(Inches(2.666), Inches(4.0), Inches(8.0), Inches(0.5))
p_ty_sub = tf_ty_sub.text_frame.paragraphs[0]
p_ty_sub.text = "We appreciate your time. Open for questions."
p_ty_sub.alignment = PP_ALIGN.CENTER
p_ty_sub.font.size = Pt(22)
p_ty_sub.font.color.rgb = ACCENT2

# ─── SAVE ────────────────────────────────────────────────────────────────────
output = "Quizzify_Presentation.pptx"
prs.save(output)
print(f"✅  Saved → {output}")
